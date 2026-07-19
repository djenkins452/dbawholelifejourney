"""
Tests for Phase 1.3 Milestone 4 — Office document perception (DOCX/XLSX/PPTX).

Deterministic text extraction via python-docx / openpyxl / python-pptx, through
the SAME perception pipeline as PDF/audio/video. Only the extractor differs.
Real Office files are generated in-test.
"""
import io

from django.contrib.auth import get_user_model
from django.test import TestCase

from apps.ai.perception import (
    _OFFICE_DOCX,
    _OFFICE_PPTX,
    _OFFICE_XLSX,
    is_perceivable,
    perceive,
)
from apps.ai.tasks import perceive_artifact
from apps.capture.models import MultimodalArtifact

User = get_user_model()


def _docx_bytes(paragraphs, table_rows=None):
    import docx
    d = docx.Document()
    for p in paragraphs:
        d.add_paragraph(p)
    if table_rows:
        t = d.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r, row in enumerate(table_rows):
            for c, val in enumerate(row):
                t.rows[r].cells[c].text = val
    buf = io.BytesIO()
    d.save(buf)
    return buf.getvalue()


def _xlsx_bytes(sheets):
    import openpyxl
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    for name, rows in sheets.items():
        ws = wb.create_sheet(title=name)
        for row in rows:
            ws.append(row)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _pptx_bytes(slides):
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for lines in slides:
        slide = prs.slides.add_slide(blank)
        tx = slide.shapes.add_textbox(0, 0, 500, 500).text_frame
        tx.text = lines[0]
        for extra in lines[1:]:
            tx.add_paragraph().text = extra
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


class OfficePerceiveTests(TestCase):
    def test_office_types_perceivable(self):
        for ct in (_OFFICE_DOCX, _OFFICE_XLSX, _OFFICE_PPTX):
            self.assertTrue(is_perceivable(ct), ct)

    def test_docx(self):
        raw = _docx_bytes(
            ["Compensation Plan", "Base salary: $120,000"],
            table_rows=[["Component", "Amount"], ["Bonus", "15%"]])
        r = perceive(_OFFICE_DOCX, raw)
        self.assertEqual(r["status"], "done")
        self.assertIn("Base salary: $120,000", r["text"])
        self.assertIn("Bonus | 15%", r["text"])        # table cell extraction

    def test_xlsx(self):
        raw = _xlsx_bytes({"Payroll": [["Name", "Salary"], ["Alex", 95000], ["Sam", 88000]]})
        r = perceive(_OFFICE_XLSX, raw)
        self.assertEqual(r["status"], "done")
        self.assertEqual(r["page_count"], 1)
        self.assertIn("[Sheet: Payroll]", r["text"])
        self.assertIn("Alex | 95000", r["text"])

    def test_pptx(self):
        raw = _pptx_bytes([["Q3 Review", "Revenue up 12%"], ["Next steps", "Ship the launch"]])
        r = perceive(_OFFICE_PPTX, raw)
        self.assertEqual(r["status"], "done")
        self.assertEqual(r["page_count"], 2)
        self.assertIn("[Slide 1]", r["text"])
        self.assertIn("Revenue up 12%", r["text"])
        self.assertIn("Ship the launch", r["text"])

    def test_corrupt_office_fails_gracefully(self):
        r = perceive(_OFFICE_DOCX, b"PK\x03\x04 not really a docx")
        self.assertIn(r["status"], ("failed", "unsupported"))
        self.assertEqual(r["text"], "")


class OfficeTaskTests(TestCase):
    def setUp(self):
        self.user = User.objects.create_user(email="office@ex.com", password="x")

    def test_task_extracts_docx(self):
        import base64
        raw = _docx_bytes(["Insurance summary", "Deductible: $2,000"])
        art = MultimodalArtifact.objects.create(
            user=self.user, sha256="a" * 64, content_type=_OFFICE_DOCX, kind="document",
            perception_status=MultimodalArtifact.PERCEPTION_PENDING)
        result = perceive_artifact(art.id, base64.b64encode(raw).decode())
        self.assertEqual(result["result"], "done")
        art.refresh_from_db()
        self.assertIn("Deductible: $2,000", art.extracted_text)
