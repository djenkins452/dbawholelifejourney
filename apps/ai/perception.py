"""
Deterministic perception — WLJ's mechanical decoders for uploaded artifacts.

The ONE place content is turned into text/transcript the conversational model can
read. This is DETERMINISTIC DECODE, not interpretation: pdfplumber returns the
characters literally present in a PDF; a future Whisper step returns the literal
transcript. WLJ never summarizes, understands, or reasons here — the model does
all perception/reasoning over the returned text
(docs/WLJ_MULTIMODAL_INTAKE_ARCHITECTURE.md §2, §3).

Dispatched by content type so every attachment reuses the SAME arrival pipeline
and only THIS component varies. Milestone 1 implements PDF; documents/audio/video
plug in the same way. Pure + defensive: never raises; returns a status.
"""
from __future__ import annotations

import io
import logging

logger = logging.getLogger(__name__)

# Cap the stored extracted text (Postgres TEXT is fine, but keep a sane ceiling
# so a pathological PDF can't store tens of MB). Full doc still fits for the vast
# majority of real documents (a 50-page policy ≈ 120K chars).
MAX_EXTRACTED_CHARS = 400_000

# Content types that have a deterministic extractor TODAY. Grows one milestone at
# a time (documents → audio → video); everything else stores + surfaces without
# perception until its extractor lands.
_OFFICE_DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_OFFICE_XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
_OFFICE_PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

PERCEIVABLE_TYPES = frozenset({
    "application/pdf",
    _OFFICE_DOCX, _OFFICE_XLSX, _OFFICE_PPTX,
    "audio/mpeg", "audio/mp4", "audio/wav", "audio/aac",
    "video/mp4", "video/quicktime",
})

# Video sampling: how many representative frames, and their downscaled longest edge.
MAX_VIDEO_FRAMES = 8
VIDEO_FRAME_SIZE = 512

# content_type → filename (extension) the shared transcription capability uses to
# tell Whisper the format.
_AUDIO_FILENAMES = {
    "audio/mpeg": "audio.mp3",
    "audio/mp4": "audio.m4a",   # our sniffer reports M4A (AAC-in-MP4) as audio/mp4
    "audio/wav": "audio.wav",
    "audio/aac": "audio.aac",   # converted to mp3 by the transcription core
}


def is_perceivable(content_type: str) -> bool:
    """True if WLJ has a deterministic extractor for this content type today."""
    return (content_type or "").lower() in PERCEIVABLE_TYPES


def perceive(content_type: str, raw: bytes) -> dict:
    """Extract deterministic text from raw bytes based on content type.

    Returns a dict: {status, text, page_count}. `status` is one of the
    MultimodalArtifact.PERCEPTION_* string values ('done' / 'failed' /
    'unsupported'). Never raises.
    """
    ct = (content_type or "").lower()
    try:
        if ct == "application/pdf":
            return _with_frames(_perceive_pdf(raw))
        if ct == _OFFICE_DOCX:
            return _with_frames(_perceive_docx(raw))
        if ct == _OFFICE_XLSX:
            return _with_frames(_perceive_xlsx(raw))
        if ct == _OFFICE_PPTX:
            return _with_frames(_perceive_pptx(raw))
        if ct in _AUDIO_FILENAMES:
            return _with_frames(_perceive_audio(raw, ct))
        if ct in ("video/mp4", "video/quicktime"):
            return _perceive_video(raw, ct)
        return {"status": "unsupported", "text": "", "page_count": None, "frames": []}
    except Exception as exc:  # pragma: no cover - defensive; perception never breaks a turn
        logger.warning("perceive failed for %s (%s)", ct, exc, exc_info=True)
        return {"status": "failed", "text": "", "page_count": None, "frames": []}


def _with_frames(result):
    """Ensure every perceive() result carries a `frames` key (only video fills it)."""
    result.setdefault("frames", [])
    return result


def _perceive_docx(raw: bytes) -> dict:
    """Extract text from a Word document (paragraphs + table cells) via python-docx."""
    import docx

    doc = docx.Document(io.BytesIO(raw))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    text = "\n".join(parts).strip()[:MAX_EXTRACTED_CHARS]
    if not text:
        return {"status": "unsupported", "text": "", "page_count": None}
    return {"status": "done", "text": text, "page_count": None}


def _perceive_xlsx(raw: bytes) -> dict:
    """Extract cell values from a spreadsheet (sheet by sheet) via openpyxl."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    parts = []
    total = 0
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            vals = [str(c) for c in row if c is not None and str(c).strip()]
            if vals:
                line = " | ".join(vals)
                parts.append(line)
                total += len(line)
        if total >= MAX_EXTRACTED_CHARS:
            break
    try:
        n_sheets = len(wb.worksheets)
    finally:
        wb.close()
    text = "\n".join(parts).strip()[:MAX_EXTRACTED_CHARS]
    if not text:
        return {"status": "unsupported", "text": "", "page_count": None}
    return {"status": "done", "text": text, "page_count": n_sheets}


def _perceive_pptx(raw: bytes) -> dict:
    """Extract slide text (per slide, in order) via python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    parts = []
    n = 0
    for i, slide in enumerate(prs.slides, start=1):
        n = i
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        lines.append(line)
        if lines:
            parts.append(f"[Slide {i}]\n" + "\n".join(lines))
        if sum(len(p) for p in parts) >= MAX_EXTRACTED_CHARS:
            break
    text = "\n\n".join(parts).strip()[:MAX_EXTRACTED_CHARS]
    if not text:
        return {"status": "unsupported", "text": "", "page_count": n or None}
    return {"status": "done", "text": text, "page_count": n or None}


def _perceive_audio(raw: bytes, content_type: str) -> dict:
    """Transcribe audio via the ONE shared transcription capability (Capture's
    Whisper integration). WLJ produces the literal transcript; the model reasons
    over it. No second transcription system.
    """
    from apps.capture.services.transcription import transcription_service

    filename = _AUDIO_FILENAMES.get(content_type, "audio.mp3")
    transcript = (transcription_service.transcribe_bytes(raw, filename) or "").strip()
    if not transcript:
        return {"status": "unsupported", "text": "", "page_count": None}
    if len(transcript) > MAX_EXTRACTED_CHARS:
        transcript = transcript[:MAX_EXTRACTED_CHARS]
    return {"status": "done", "text": transcript, "page_count": None}


def _perceive_video(raw: bytes, content_type: str) -> dict:
    """Video perception = DUAL deterministic decode, reusing existing capabilities:
      • representative FRAMES (ffmpeg) → delivered to the model's image-perception
        path so it can SEE motion/form/scene (squat form, golf swing, "what am I
        doing");
      • the AUDIO TRACK transcript via the ONE shared transcription capability
        (meetings, narration).
    WLJ only decodes (samples frames + transcribes); the model reasons. Both are
    best-effort; either alone yields a usable result.
    """
    ext = "mov" if content_type == "video/quicktime" else "mp4"

    try:
        frames = _extract_frames(raw, ext)
    except Exception as exc:
        logger.warning("video frame extraction failed (%s)", exc)
        frames = []

    transcript = ""
    try:
        from apps.capture.services.transcription import transcription_service
        transcript = (transcription_service.transcribe_bytes(raw, f"video.{ext}") or "").strip()
    except Exception as exc:
        logger.info("video audio-track transcription skipped (%s)", exc)
        transcript = ""

    parts = []
    if frames:
        ts = ", ".join(f"{f['t']:.1f}s" for f in frames)
        parts.append(
            f"[Video] {len(frames)} representative frames were sampled at {ts} and "
            "provided to you AS IMAGES this turn — view them in order to evaluate the "
            "video (motion, form, scene, what is happening)."
        )
    if transcript:
        parts.append(f"[Audio transcript]\n{transcript}")

    if not frames and not transcript:
        return {"status": "unsupported", "text": "", "page_count": None, "frames": []}

    text = "\n\n".join(parts)
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS]
    return {"status": "done", "text": text,
            "page_count": len(frames) or None, "frames": frames}


def _extract_frames(raw: bytes, ext: str,
                    max_frames: int = MAX_VIDEO_FRAMES, size: int = VIDEO_FRAME_SIZE):
    """Sample up to `max_frames` evenly-spaced, downscaled JPEG frames via ffmpeg.
    Returns [{"t": seconds, "b64": jpeg_base64}]. Requires ffmpeg/ffprobe; returns
    [] (never raises to the caller's detriment) if unavailable.
    """
    import base64
    import glob
    import os
    import shutil
    import subprocess
    import tempfile

    inp = None
    outdir = tempfile.mkdtemp(prefix="wlj-frames-")
    frames = []
    try:
        with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as f:
            f.write(raw)
            inp = f.name

        # Duration (ffprobe) → evenly-centered sample timestamps.
        duration = None
        try:
            r = subprocess.run(
                ["ffprobe", "-v", "error", "-show_entries", "format=duration",
                 "-of", "default=noprint_wrappers=1:nokey=1", inp],
                capture_output=True, timeout=60,
            )
            duration = float((r.stdout or b"").decode().strip())
        except Exception:
            duration = None

        vf = f"scale={size}:-2"
        if duration and duration > 0:
            times = [round(duration * (i + 0.5) / max_frames, 2) for i in range(max_frames)]
            for i, t in enumerate(times):
                out = os.path.join(outdir, f"f{i:03d}.jpg")
                subprocess.run(
                    ["ffmpeg", "-y", "-ss", str(t), "-i", inp,
                     "-frames:v", "1", "-vf", vf, "-q:v", "5", out],
                    capture_output=True, timeout=60,
                )
                if os.path.exists(out):
                    with open(out, "rb") as fh:
                        frames.append({"t": t, "b64": base64.b64encode(fh.read()).decode()})
        else:
            # No duration — grab up to max_frames at 1 fps.
            subprocess.run(
                ["ffmpeg", "-y", "-i", inp, "-vf", f"fps=1,{vf}",
                 "-frames:v", str(max_frames), "-q:v", "5",
                 os.path.join(outdir, "f%03d.jpg")],
                capture_output=True, timeout=120,
            )
            for i, p in enumerate(sorted(glob.glob(os.path.join(outdir, "*.jpg")))):
                with open(p, "rb") as fh:
                    frames.append({"t": float(i), "b64": base64.b64encode(fh.read()).decode()})
        return frames
    finally:
        if inp:
            try:
                os.unlink(inp)
            except OSError:
                pass
        try:
            shutil.rmtree(outdir)
        except OSError:
            pass


def _perceive_pdf(raw: bytes) -> dict:
    """Extract text from a text-based PDF via pdfplumber (page-delimited)."""
    import pdfplumber

    pages_text = []
    with pdfplumber.open(io.BytesIO(raw)) as pdf:
        page_count = len(pdf.pages)
        for idx, page in enumerate(pdf.pages, start=1):
            txt = page.extract_text() or ""
            if txt.strip():
                pages_text.append(f"[Page {idx}]\n{txt.strip()}")
            # Bound total work — stop once we've gathered enough characters.
            if sum(len(p) for p in pages_text) >= MAX_EXTRACTED_CHARS:
                break

    text = "\n\n".join(pages_text).strip()
    if len(text) > MAX_EXTRACTED_CHARS:
        text = text[:MAX_EXTRACTED_CHARS]

    if not text:
        # A scanned/image-only PDF yields no text here — OCR is a later step.
        return {"status": "unsupported", "text": "", "page_count": page_count}
    return {"status": "done", "text": text, "page_count": page_count}
